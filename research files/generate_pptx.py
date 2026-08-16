import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("python-pptx is not installed.")
    sys.exit(1)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette - Professional Corporate Slate Theme
    COLOR_BG = RGBColor(15, 23, 42)        # Slate 900 #0F172A
    COLOR_CARD = RGBColor(30, 41, 59)      # Slate 800 #1E293B
    COLOR_BORDER = RGBColor(51, 65, 85)    # Slate 700 #334155
    COLOR_ACCENT = RGBColor(14, 165, 233)  # Sky 500 #0EA5E9
    COLOR_PURPLE = RGBColor(139, 92, 246)  # Violet 500 #8B5CF6
    COLOR_GREEN = RGBColor(34, 197, 94)    # Emerald 500 #22C55E
    COLOR_AMBER = RGBColor(245, 158, 11)   # Amber 500 #F59E0B
    COLOR_TEXT = RGBColor(248, 250, 252)   # White #F8FAFC
    COLOR_MUTED = RGBColor(148, 163, 184)  # Slate 400 #94A3B8

    def add_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.color.rgb = COLOR_BG

    def add_header(slide, title_text, category_text="FORGEQA PRODUCT OVERVIEW"):
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.45), Inches(3.2), Inches(0.35))
        badge.fill.solid()
        badge.fill.fore_color.rgb = COLOR_CARD
        badge.line.color.rgb = COLOR_ACCENT
        tf_b = badge.text_frame
        p_b = tf_b.paragraphs[0]
        p_b.text = category_text.upper()
        p_b.font.size = Pt(10)
        p_b.font.bold = True
        p_b.font.color.rgb = COLOR_ACCENT
        p_b.alignment = PP_ALIGN.CENTER

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(10.5), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT

        logo_box = slide.shapes.add_textbox(Inches(10.5), Inches(0.4), Inches(2.0), Inches(0.5))
        tf_l = logo_box.text_frame
        p_l = tf_l.paragraphs[0]
        p_l.text = "ForgeQA.in"
        p_l.font.size = Pt(16)
        p_l.font.bold = True
        p_l.font.color.rgb = COLOR_ACCENT
        p_l.alignment = PP_ALIGN.RIGHT

    def add_footer(slide, current_slide, total_slides=8):
        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.4))
        tf = footer_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"ForgeQA Product Presentation  |  Built Features Overview  |  Slide {current_slide} of {total_slides}"
        p.font.size = Pt(9)
        p.font.color.rgb = COLOR_MUTED

    blank_layout = prs.slide_layouts[6]

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    add_bg(slide1)

    dec = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(0.15), Inches(3.8))
    dec.fill.solid()
    dec.fill.fore_color.rgb = COLOR_ACCENT
    dec.line.color.rgb = COLOR_ACCENT

    t_box = slide1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11.0), Inches(3.8))
    tf1 = t_box.text_frame
    tf1.word_wrap = True

    p1 = tf1.paragraphs[0]
    p1.text = "ForgeQA"
    p1.font.size = Pt(54)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TEXT

    p2 = tf1.add_paragraph()
    p2.text = "AI-Powered Test Generation & Automation Platform"
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_ACCENT
    p2.space_before = Pt(8)

    p3 = tf1.add_paragraph()
    p3.text = "From Natural Language Requirements to Executable Test Scripts, RAG Knowledge Base & Regression Runs"
    p3.font.size = Pt(14)
    p3.font.color.rgb = COLOR_MUTED
    p3.space_before = Pt(14)

    # Simple info card
    p_card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.8), Inches(11.733), Inches(1.1))
    p_card.fill.solid()
    p_card.fill.fore_color.rgb = COLOR_CARD
    p_card.line.color.rgb = COLOR_BORDER
    tf_pc = p_card.text_frame
    tf_pc.word_wrap = True
    p_pc1 = tf_pc.paragraphs[0]
    p_pc1.text = "Built Platform Capabilities Overview  |  Stack: React 19, Node.js, Express, MongoDB, Multi-AI Providers"
    p_pc1.font.size = Pt(12)
    p_pc1.font.bold = True
    p_pc1.font.color.rgb = COLOR_TEXT
    p_pc2 = tf_pc.add_paragraph()
    p_pc2.text = "Website: https://forgeqa.in  |  Application Engine Version: 0.1.0"
    p_pc2.font.size = Pt(11)
    p_pc2.font.color.rgb = COLOR_MUTED
    p_pc2.space_before = Pt(4)


    # -------------------------------------------------------------
    # SLIDE 2: What is ForgeQA (Core Capabilities)
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    add_bg(slide2)
    add_header(slide2, "Core Capabilities Built in ForgeQA")
    add_footer(slide2, 2)

    features = [
        ("AI Test Case Generator", "Converts plain text requirements into categorized test matrices (Positive, Negative, Validation, Edge cases).", COLOR_ACCENT),
        ("Multi-Framework Script Generator", "Generates executable code in Playwright, Cypress, Selenium, Puppeteer (TS, JS, Python, Java, C#).", COLOR_PURPLE),
        ("BYOK Multi-AI Provider Engine", "Connect your own API key for Gemini, OpenAI, Claude, Groq, OpenRouter, or OpenCode.", COLOR_GREEN),
        ("RAG Knowledge Base", "Upload PDF, Word (DOCX), TXT, or SharePoint docs to ground test generation in company specs.", COLOR_AMBER),
        ("Regression & Build Engine", "Upload Web & Android APK builds, configure webhooks, and trigger automated regression runs.", COLOR_ACCENT),
        ("Suites & License Governance", "Group test cases into color-coded suites, assign team roles, and secure data with Screen Capture Shield.", COLOR_PURPLE)
    ]

    for idx, (title, desc, color) in enumerate(features):
        col = idx % 3
        row = idx // 3
        left = Inches(0.8 + col * 3.95)
        top = Inches(1.8 + row * 2.5)

        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.8), Inches(2.3))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = COLOR_BORDER

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = color

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = COLOR_TEXT
        p_d.space_before = Pt(8)


    # -------------------------------------------------------------
    # SLIDE 3: AI Test Case Matrix Generator
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    add_bg(slide3)
    add_header(slide3, "AI Test Matrix Generator (/generator)")
    add_footer(slide3, 3)

    left_card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.7), Inches(4.9))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = COLOR_CARD
    left_card.line.color.rgb = COLOR_ACCENT

    tf_l = left_card.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = Inches(0.3)
    tf_l.margin_top = Inches(0.3)

    p = tf_l.paragraphs[0]
    p.text = "Categorized Test Case Matrix"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    tc_points = [
        "Positive Test Cases: Verifies expected standard user flows.",
        "Negative Test Cases: Checks invalid inputs & failure handling.",
        "Validation Checks: Validates input field bounds & formats.",
        "Edge Cases: Identifies boundary conditions & unusual scenarios.",
        "Structured Schema: Unique ID (TC_001...), Summary, Test Steps, Expected Result, Status (Draft / Reviewed / Approved)."
    ]
    for tcp in tc_points:
        p_tp = tf_l.add_paragraph()
        p_tp.text = f"•  {tcp}"
        p_tp.font.size = Pt(11)
        p_tp.font.color.rgb = COLOR_TEXT
        p_tp.space_before = Pt(12)

    right_card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.9))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = COLOR_CARD
    right_card.line.color.rgb = COLOR_GREEN

    tf_r = right_card.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = Inches(0.3)
    tf_r.margin_top = Inches(0.3)

    p_r = tf_r.paragraphs[0]
    p_r.text = "One-Click Document Export & Storage"
    p_r.font.size = Pt(18)
    p_r.font.bold = True
    p_r.font.color.rgb = COLOR_GREEN

    export_points = [
        "Export to Excel (.xlsx): Spreadsheets formatted for QA teams.",
        "Export to PDF (.pdf): Printable test specifications.",
        "Export to Word (.docx): Document reports for documentation.",
        "Persistent History: Saves up to 50 generation runs per user.",
        "Quick Reload: Instant retrieval of past test matrices."
    ]
    for ep in export_points:
        p_ep = tf_r.add_paragraph()
        p_ep.text = f"✓  {ep}"
        p_ep.font.size = Pt(11)
        p_ep.font.color.rgb = COLOR_TEXT
        p_ep.space_before = Pt(14)


    # -------------------------------------------------------------
    # SLIDE 4: Multi-Framework Script Generator & BYOK
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    add_bg(slide4)
    add_header(slide4, "Automated Script Generator & Multi-AI BYOK")
    add_footer(slide4, 4)

    # 2 Column Cards
    sc1 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.7), Inches(4.9))
    sc1.fill.solid()
    sc1.fill.fore_color.rgb = COLOR_CARD
    sc1.line.color.rgb = COLOR_PURPLE

    tf1 = sc1.text_frame
    tf1.word_wrap = True
    tf1.margin_left = Inches(0.3)
    tf1.margin_top = Inches(0.3)

    p = tf1.paragraphs[0]
    p.text = "Frameworks & Programming Languages"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_PURPLE

    fw_points = [
        "Automation Frameworks: Playwright, Cypress, Selenium, Puppeteer",
        "Supported Languages: TypeScript, JavaScript, Python, Java, C#",
        "Target URL Configuration: Generates page object selectors for your exact app URL",
        "Options & Viewport: Headless mode toggle & configurable viewport width/height"
    ]
    for fp in fw_points:
        p_fp = tf1.add_paragraph()
        p_fp.text = f"•  {fp}"
        p_fp.font.size = Pt(12)
        p_fp.font.color.rgb = COLOR_TEXT
        p_fp.space_before = Pt(14)

    sc2 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.9))
    sc2.fill.solid()
    sc2.fill.fore_color.rgb = COLOR_CARD
    sc2.line.color.rgb = COLOR_ACCENT

    tf2 = sc2.text_frame
    tf2.word_wrap = True
    tf2.margin_left = Inches(0.3)
    tf2.margin_top = Inches(0.3)

    p_a = tf2.paragraphs[0]
    p_a.text = "Bring Your Own Key (BYOK) Providers"
    p_a.font.size = Pt(18)
    p_a.font.bold = True
    p_a.font.color.rgb = COLOR_ACCENT

    ai_providers = [
        "Google Gemini: Gemini 2.0 Flash & Gemini 2.5 Flash",
        "OpenAI: GPT-4o & GPT models",
        "Anthropic Claude: Claude 3.5 Sonnet",
        "Groq & OpenRouter: High-speed inference models",
        "OpenCode / Local AI: Custom API endpoint integration",
        "Secure Storage: Per-user API key encryption in settings"
    ]
    for ap in ai_providers:
        p_ap = tf2.add_paragraph()
        p_ap.text = f"✓  {ap}"
        p_ap.font.size = Pt(11)
        p_ap.font.color.rgb = COLOR_TEXT
        p_ap.space_before = Pt(10)


    # -------------------------------------------------------------
    # SLIDE 5: RAG Knowledge Base Engine
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    add_bg(slide5)
    add_header(slide5, "RAG Knowledge Base & Document Context (/knowledge)")
    add_footer(slide5, 5)

    kb_cards = [
        ("Supported File Formats", "Upload PDF (.pdf), Word (.docx), Plain Text (.txt), Markdown (.md), and SharePoint synced documents.", COLOR_ACCENT),
        ("Processing & Chunking", "Automatic text extraction & semantic chunking status tracking: Processing → Needs Chunking → Ready.", COLOR_PURPLE),
        ("Context-Aware AI Generation", "When generating tests, AI retrieves top relevant knowledge chunks score-ranked against requirements.", COLOR_GREEN),
        ("Doc Parsing Engine", "Integrated server parsers: pdf-parse, mammoth (Word), tesseract.js OCR, and MongoDB document storage.", COLOR_AMBER)
    ]

    for idx, (title, desc, color) in enumerate(kb_cards):
        col = idx % 2
        row = idx // 2
        left = Inches(0.8 + col * 5.95)
        top = Inches(1.8 + row * 2.5)

        card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.7), Inches(2.3))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = color

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)
        tf.margin_right = Inches(0.3)
        tf.margin_top = Inches(0.25)

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = COLOR_TEXT
        p_d.space_before = Pt(10)


    # -------------------------------------------------------------
    # SLIDE 6: Automated Regression Testing Engine
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    add_bg(slide6)
    add_header(slide6, "Automated Regression Testing Engine (/regression)")
    add_footer(slide6, 6)

    # 3 Column Cards
    reg_items = [
        ("Build Artifact Uploads", [
            "Upload Web build bundles (.zip)",
            "Upload Android application builds (.apk)",
            "Track version numbers & build sizes",
            "Store build history in system storage"
        ], COLOR_ACCENT),
        ("Webhook Integration", [
            "Configurable build webhook endpoints",
            "Auto-trigger regression runs on build updates",
            "Web & Mobile platform selectors",
            "Seamless integration into build scripts"
        ], COLOR_PURPLE),
        ("Execution Status & Artifacts", [
            "Track status: Pending, Running, Passed, Failed",
            "Detailed actual output & error log messages",
            "Screenshot artifact capture on test failures",
            "Run completion timestamps & build versioning"
        ], COLOR_GREEN)
    ]

    for i, (header, items, color) in enumerate(reg_items):
        left = Inches(0.8 + i * 3.95)
        cc = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(3.8), Inches(4.9))
        cc.fill.solid()
        cc.fill.fore_color.rgb = COLOR_CARD
        cc.line.color.rgb = color

        tf = cc.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.25)

        p = tf.paragraphs[0]
        p.text = header
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color

        for it in items:
            p_it = tf.add_paragraph()
            p_it.text = f"•  {it}"
            p_it.font.size = Pt(12)
            p_it.font.color.rgb = COLOR_TEXT
            p_it.space_before = Pt(14)


    # -------------------------------------------------------------
    # SLIDE 7: Test Suites, Security & Admin Control
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    add_bg(slide7)
    add_header(slide7, "Test Suites, Security & Governance")
    add_footer(slide7, 7)

    sec_items = [
        ("Test Suite Management (/suites)", [
            "Group test cases into custom test suites",
            "Color-coded suite badges & descriptions",
            "Track case status: Draft, Reviewed, Approved"
        ], COLOR_ACCENT),
        ("Team & License Control (/settings)", [
            "Product Key activation & validation",
            "Team roles: Admin, Member, Viewer",
            "User profile management & password reset"
        ], COLOR_PURPLE),
        ("Built-in Screen Capture Shield", [
            "Blocks PrintScreen & Win+Shift+S snipping",
            "Blocks right-click context menu & image dragging",
            "Logs security events to protect corporate IP"
        ], COLOR_AMBER)
    ]

    for i, (header, items, color) in enumerate(sec_items):
        left = Inches(0.8 + i * 3.95)
        cc = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(3.8), Inches(4.9))
        cc.fill.solid()
        cc.fill.fore_color.rgb = COLOR_CARD
        cc.line.color.rgb = color

        tf = cc.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.25)

        p = tf.paragraphs[0]
        p.text = header
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color

        for it in items:
            p_it = tf.add_paragraph()
            p_it.text = f"•  {it}"
            p_it.font.size = Pt(12)
            p_it.font.color.rgb = COLOR_TEXT
            p_it.space_before = Pt(14)


    # -------------------------------------------------------------
    # SLIDE 8: Technology Summary & Live Platform Access
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    add_bg(slide8)
    add_header(slide8, "Platform Architecture & Live Demo")
    add_footer(slide8, 8)

    top_card = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.733), Inches(2.6))
    top_card.fill.solid()
    top_card.fill.fore_color.rgb = COLOR_CARD
    top_card.line.color.rgb = COLOR_ACCENT

    tf_t = top_card.text_frame
    tf_t.word_wrap = True
    tf_t.margin_left = Inches(0.3)
    tf_t.margin_top = Inches(0.25)

    p_tt = tf_t.paragraphs[0]
    p_tt.text = "Technical Architecture & Engine Stack"
    p_tt.font.size = Pt(18)
    p_tt.font.bold = True
    p_tt.font.color.rgb = COLOR_ACCENT

    tech_stack = [
        "Frontend Application: React 19, TypeScript, Vite, TailwindCSS v4, Zustand, Framer Motion, Three.js, Lucide Icons.",
        "Backend REST API: Node.js (v24), Express, MongoDB driver, JWT authentication with session refresh cookie.",
        "Document & Parser Libraries: ExcelJS / XLSX, pdf-parse, Mammoth (docx), Tesseract.js (OCR), Multer upload handling."
    ]
    for ts in tech_stack:
        p_ts = tf_t.add_paragraph()
        p_ts.text = f"•  {ts}"
        p_ts.font.size = Pt(11)
        p_ts.font.color.rgb = COLOR_TEXT
        p_ts.space_before = Pt(8)

    bottom_card = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.7), Inches(11.733), Inches(2.0))
    bottom_card.fill.solid()
    bottom_card.fill.fore_color.rgb = COLOR_CARD
    bottom_card.line.color.rgb = COLOR_GREEN

    tf_b = bottom_card.text_frame
    tf_b.word_wrap = True
    tf_b.margin_left = Inches(0.3)
    tf_b.margin_top = Inches(0.2)

    p_bt = tf_b.paragraphs[0]
    p_bt.text = "Experience ForgeQA Live"
    p_bt.font.size = Pt(18)
    p_bt.font.bold = True
    p_bt.font.color.rgb = COLOR_TEXT

    p_bi = tf_b.add_paragraph()
    p_bi.text = "Website: https://forgeqa.in  |  API Server: node server/index.js"
    p_bi.font.size = Pt(14)
    p_bi.font.color.rgb = COLOR_GREEN
    p_bi.space_before = Pt(8)

    p_bs = tf_b.add_paragraph()
    p_bs.text = "Thank you! Open for Questions & Demonstration."
    p_bs.font.size = Pt(12)
    p_bs.font.color.rgb = COLOR_MUTED
    p_bs.space_before = Pt(8)

    output_path = os.path.join("d:\\forgeqa.in", "ForgeQA_Client_Presentation.pptx")
    prs.save(output_path)
    print(f"Successfully generated real-data PowerPoint presentation at: {output_path}")

if __name__ == "__main__":
    create_presentation()
